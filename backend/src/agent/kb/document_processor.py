"""文档处理流水线：解析→分块→向量化→写入 Milvus。

支持格式：PDF/Word/Excel/PPT/Markdown/TXT
状态流转：pending → parsing → vectorizing → done / failed
"""

import os
from pathlib import Path
from typing import Optional

from loguru import logger


def parse_pdf(file_path: str) -> str:
    """解析 PDF 文件"""
    from PyPDF2 import PdfReader
    reader = PdfReader(file_path)
    text_parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            text_parts.append(text)
    return "\n\n".join(text_parts)


def parse_docx(file_path: str) -> str:
    """解析 Word 文件"""
    from docx import Document
    doc = Document(file_path)
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                text_parts.append(row_text)
    return "\n\n".join(text_parts)


def parse_excel(file_path: str) -> str:
    """解析 Excel 文件"""
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    text_parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        text_parts.append(f"[Sheet: {sheet_name}]")
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                text_parts.append(row_text)
    wb.close()
    return "\n\n".join(text_parts)


def parse_pptx(file_path: str) -> str:
    """解析 PPT 文件"""
    from pptx import Presentation
    prs = Presentation(file_path)
    text_parts = []
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_texts = [f"[Slide {slide_num}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text)
        if len(slide_texts) > 1:
            text_parts.append("\n".join(slide_texts))
    return "\n\n".join(text_parts)


def parse_text(file_path: str) -> str:
    """解析纯文本文件（Markdown/TXT）"""
    with open(file_path, encoding="utf-8", errors="ignore") as f:
        return f.read()


PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".doc": parse_docx,
    ".xlsx": parse_excel,
    ".xls": parse_excel,
    ".pptx": parse_pptx,
    ".ppt": parse_pptx,
    ".md": parse_text,
    ".txt": parse_text,
    ".csv": parse_text,
}

SUPPORTED_EXTENSIONS = set(PARSERS.keys())


def parse_document(file_path: str) -> str:
    """根据文件扩展名选择解析器，返回文本内容"""
    ext = Path(file_path).suffix.lower()
    parser = PARSERS.get(ext)
    if not parser:
        raise ValueError(f"不支持的文件格式: {ext}")
    return parser(file_path)


def chunk_text(text: str, chunk_size: int = 512, chunk_overlap: int = 64) -> list[str]:
    """将文本分块，每块约 chunk_size 字符，重叠 chunk_overlap 字符"""
    if not text or not text.strip():
        return []

    chunks = []
    start = 0
    text_len = len(text)

    while start < text_len:
        end = start + chunk_size
        chunk = text[start:end]

        if chunk.strip():
            chunks.append(chunk.strip())

        start = end - chunk_overlap
        if start >= text_len:
            break

    return chunks


def process_uploaded_file(
    file_path: str,
    file_id: int,
    user_id: int,
    title: str,
    embedding_fn=None,
    milvus_collection=None,
) -> dict:
    """完整的文件处理流水线：解析→分块→向量化→写入 Milvus

    Args:
        file_path: 文件本地路径
        file_id: 数据库中的文件 ID
        user_id: 用户 ID
        title: 文件标题
        embedding_fn: 嵌入函数，接受 texts 返回 embeddings
        milvus_collection: Milvus collection 实例

    Returns:
        {"status": "done"|"failed", "chunk_count": int, "error": str|None}
    """
    try:
        # 1. 解析
        logger.info(f"[DocProcessor] 解析文件: {title}")
        text = parse_document(file_path)

        if not text or not text.strip():
            return {"status": "failed", "chunk_count": 0, "error": "文件内容为空"}

        # 2. 分块
        chunks = chunk_text(text)
        if not chunks:
            return {"status": "failed", "chunk_count": 0, "error": "分块后无有效内容"}

        logger.info(f"[DocProcessor] 分块完成: {len(chunks)} chunks")

        # 3. 向量化并写入 Milvus
        if embedding_fn and milvus_collection:
            embeddings = embedding_fn(chunks)
            ids = [f"kb_{user_id}_{file_id}_{i}" for i in range(len(chunks))]
            metadatas = [
                {
                    "user_id": user_id,
                    "file_id": file_id,
                    "chunk_index": i,
                    "source_type": "upload",
                    "title": title,
                    "file_ext": Path(file_path).suffix.lower(),
                }
                for i in range(len(chunks))
            ]

            milvus_collection.upsert(
                ids=ids,
                embeddings=embeddings,
                documents=chunks,
                metadatas=metadatas,
            )
            logger.info(f"[DocProcessor] 向量写入完成: {len(chunks)} 条")

        return {"status": "done", "chunk_count": len(chunks), "error": None}

    except Exception as e:
        logger.error(f"[DocProcessor] 处理失败: {e}")
        return {"status": "failed", "chunk_count": 0, "error": str(e)}
